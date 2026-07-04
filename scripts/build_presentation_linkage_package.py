#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUT_DIR = ROOT / "evaluation" / "presentation-linkage"
FONT = "Malgun Gothic"

COLORS = {
    "ink": "1F2937",
    "muted": "64748B",
    "line": "CBD5E1",
    "paper": "F8FAFC",
    "white": "FFFFFF",
    "teal": "0F766E",
    "cyan": "0891B2",
    "blue": "2563EB",
    "indigo": "4F46E5",
    "amber": "B45309",
    "rose": "BE123C",
    "green": "15803D",
}


def read_json(relative_path: str) -> dict[str, Any]:
    path = ROOT / relative_path
    with path.open("r", encoding="utf-8") as handle:
        value = json.load(handle)
    if not isinstance(value, dict):
        raise ValueError(f"Expected JSON object: {path}")
    return value


def write_json(path: Path, value: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def write_text(path: Path, value: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(value, encoding="utf-8")


def hex_color(value: str) -> RGBColor:
    value = value.strip().lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds").replace("+00:00", "Z")


def graph_counts(graph: dict[str, Any]) -> dict[str, int]:
    elements = graph.get("elements", {})
    nodes = elements.get("nodes", [])
    edges = elements.get("edges", [])
    return {"nodes": len(nodes), "edges": len(edges)}


def extract_pptx_slide_text(path: Path) -> list[dict[str, Any]]:
    if not path.exists():
        return []
    prs = Presentation(path)
    slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text = str(shape.text or "").strip()
                if text:
                    texts.append(text.replace("\n", " | "))
        slides.append({"slide": slide_index, "text": " || ".join(texts)})
    return slides


def agenda_deck_cross_check() -> dict[str, Any]:
    relative_path = Path("10_작업산출물") / "7.4_발표덱" / "2026기후시민회의_숙의의제.pptx"
    deck_path = ROOT.parent / relative_path
    slides = extract_pptx_slide_text(deck_path)
    signals = [
        item
        for item in slides
        if item["slide"] in {13, 14, 15, 18}
        or any(term in item["text"] for term in ["1위 의제", "공동 2위 의제", "우리가 발견한 것"])
    ]
    return {
        "path": relative_path.as_posix(),
        "exists": deck_path.exists(),
        "slides": len(slides),
        "signals": signals,
        "current_three_agendas": [
            {
                "slot": "1위",
                "conclusion": "전 생애주기 탄소중립 교육 체계 구축",
                "source_decision_id": "agenda-education-citizen-participation",
                "deck_slide": 13,
                "status": "현재 숙의의제 발표덱에 별도 의제로 제시",
            },
            {
                "slot": "공동 2위",
                "conclusion": "시민의식 개선 · 참여 활성화",
                "source_decision_id": "agenda-citizen-participation-original",
                "deck_slide": 14,
                "status": "현재 숙의의제 발표덱에 별도 의제로 제시",
            },
            {
                "slot": "공동 2위",
                "conclusion": "자원순환 · 생활폐기물 감축",
                "source_decision_id": "agenda-resource-circulation",
                "deck_slide": 15,
                "status": "현재 숙의의제 발표덱에 별도 의제로 제시",
            },
        ],
    }


def short_vote(decision: dict[str, Any]) -> str:
    vote = decision.get("vote")
    if vote:
        return str(vote)
    rank = decision.get("vote_rank")
    if rank is None:
        return "확정 근거 부족"
    return f"{rank}위"


def build_evidence_map() -> dict[str, Any]:
    input_coverage = read_json("evaluation/input-coverage/input-coverage-report.json")
    final_report = read_json("evaluation/ontology-final-decisions/final-decision-ontology-report.json")
    sources = read_json("public/workshop-graph/sources.json")
    source_coverage = read_json("public/workshop-graph/data/source-coverage-2026-06-13.json")
    process_graph = read_json("public/workshop-graph/data/final-process-to-conclusion-0704.json")
    final_regulation = read_json("public/workshop-graph/data/final-regulation-decisions-0704.json")
    final_agenda = read_json("public/workshop-graph/data/final-agenda-decisions-0704.json")

    coverage_inventory = input_coverage["inventory"]
    process_coverage = final_report["coverage"]["process"]
    regulation_items = final_report["coverage"]["regulation"]
    agenda_items = final_report["coverage"]["agenda"]

    agenda_by_id = {item["id"]: item for item in agenda_items}
    deck_check = agenda_deck_cross_check()
    final_slots = []
    for deck_item in deck_check["current_three_agendas"]:
        decision = agenda_by_id[deck_item["source_decision_id"]]
        if deck_item["source_decision_id"] == "agenda-education-citizen-participation":
            linkage = [
                "사전 접수와 조별 숙의에서 교육 필요성이 반복됨",
                "투표에서 가장 앞선 후보로 확인됨",
                "현재 숙의의제 발표덱 13쪽에서 별도 의제로 설명됨",
            ]
        elif deck_item["source_decision_id"] == "agenda-citizen-participation-original":
            linkage = [
                "시민의식 개선과 참여 활성화 맥락이 교육·거버넌스 논의와 연결됨",
                "투표에서 공동 2위권 후보로 확인됨",
                "현재 숙의의제 발표덱 14쪽에서 별도 의제로 설명됨",
            ]
        else:
            linkage = [
                "A조 자원순환형 배달 문화와 B조 생활폐기물 감축 논의가 만남",
                "투표에서 공동 2위권 후보로 확인됨",
                "현재 숙의의제 발표덱 15쪽에서 별도 의제로 설명됨",
            ]
        final_slots.append(
            {
                "slot": deck_item["slot"],
                "status": deck_item["status"],
                "conclusion": deck_item["conclusion"],
                "deck_slide": deck_item["deck_slide"],
                "source_decisions": [decision],
                "linkage": linkage,
            }
        )

    scenario_variant_slots = [
        {
            "slot": "적응",
            "status": "v6 시나리오상 통합 설명",
            "conclusion": "생애주기 탄소중립 교육 및 시민의식 개선",
            "source_decisions": [
                agenda_by_id["agenda-education-citizen-participation"],
                agenda_by_id["agenda-citizen-participation-original"],
            ],
            "linkage": [
                "사전 접수와 조별 숙의에서 교육·시민참여 맥락이 반복됨",
                "교육 의제가 평가 1위, 시민참여 의제가 공동 2위로 확인됨",
                "발표 시나리오에서 시민참여를 교육 의제에 통합해 적응 분과 의제로 설명함",
            ],
        },
        {
            "slot": "감축1",
            "status": "v6 시나리오상 확정 설명 가능",
            "conclusion": "자원순환·생활폐기물 감축",
            "source_decisions": [agenda_by_id["agenda-resource-circulation"]],
            "linkage": [
                "사전 접수 분석에서 소비·생활양식과 자원순환 흐름이 크게 나타남",
                "A조 자원순환형 배달 문화와 B조 생활폐기물 감축 논의를 통합함",
                "평가 결과 공동 2위권으로 확인되어 감축1 분과 의제로 설명함",
            ],
        },
        {
            "slot": "감축2",
            "status": "v6 시나리오상 확정명 증거 부족",
            "conclusion": "새로운 의제 슬롯",
            "source_decisions": [],
            "linkage": [
                "시민의식 개선 후보가 교육 의제에 통합되며 세 번째 자리가 열림",
                "발표 시나리오 7쪽·10쪽에 7.4 당일 확정 후 수정 표시가 남아 있음",
                "현재 repo 증거만으로는 별도 최종 의제명을 단정할 수 없음",
            ],
        },
    ]

    public_source_ids = [source["id"] for source in sources["sources"]]
    return {
        "generated_at": now_iso(),
        "summary": {
            "a_only_public_menu": False,
            "all_data_completely_reflected": False,
            "linkage_supported": True,
            "recommended_framing": "숙의가 기준과 선택지를 만들고, 투표와 통합 판단이 결론을 고정했다.",
        },
        "public_menu": {
            "default": sources["default"],
            "categories": sources["categories"],
            "source_ids": public_source_ids,
            "has_live_a_t1": "live-A_t1" in public_source_ids,
        },
        "coverage": {
            "source_coverage_graph": graph_counts(source_coverage),
            "process_graph": graph_counts(process_graph),
            "final_regulation_graph": graph_counts(final_regulation),
            "final_agenda_graph": graph_counts(final_agenda),
            "original_files": coverage_inventory["sourceCoverage"]["originalFiles"],
            "workflow_markdown_files": len(coverage_inventory["workflowInputMarkdownFiles"]),
            "ready_sessions": process_coverage["ready_sessions"],
            "partial_or_review_sessions": process_coverage["partial_or_review_sessions"],
            "partial_or_review_labels": process_coverage["partial_or_review_labels"],
            "input_gap_items": input_coverage["coverage"]["partialOrGapSessions"],
        },
        "regulation_decisions": regulation_items,
        "agenda_candidates": agenda_items,
        "agenda_deck_cross_check": deck_check,
        "agenda_final_slots": final_slots,
        "scenario_variant_slots": scenario_variant_slots,
        "caveats": final_report["caveats"],
        "source_files": {
            "input_coverage": "evaluation/input-coverage/input-coverage-report.json",
            "final_decision_report": "evaluation/ontology-final-decisions/final-decision-ontology-report.json",
            "sources": "public/workshop-graph/sources.json",
            "source_coverage_graph": "public/workshop-graph/data/source-coverage-2026-06-13.json",
            "process_graph": "public/workshop-graph/data/final-process-to-conclusion-0704.json",
            "final_regulation_graph": "public/workshop-graph/data/final-regulation-decisions-0704.json",
            "final_agenda_graph": "public/workshop-graph/data/final-agenda-decisions-0704.json",
            "current_agenda_deck": deck_check["path"],
            "regulation_scenario": "10_작업산출물/7.4_발표덱/운영규정_v6/발표시나리오_운영규정_이대진.md",
            "agenda_scenario": "10_작업산출물/7.4_발표덱/의제결과_v6/발표시나리오_의제선정결과_김영현.md",
        },
    }


def set_text_style(text_frame: Any, font_size: int, color: str, bold: bool = False) -> None:
    for paragraph in text_frame.paragraphs:
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = hex_color(color)


def add_text(
    slide: Any,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    font_size: int = 18,
    color: str = COLORS["ink"],
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    valign: MSO_ANCHOR = MSO_ANCHOR.TOP,
) -> Any:
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = Inches(0.08)
    frame.margin_right = Inches(0.08)
    frame.margin_top = Inches(0.04)
    frame.margin_bottom = Inches(0.04)
    frame.vertical_anchor = valign
    first = True
    for line in text.split("\n"):
        paragraph = frame.paragraphs[0] if first else frame.add_paragraph()
        paragraph.text = line
        paragraph.alignment = align
        paragraph.font.name = FONT
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        paragraph.font.color.rgb = hex_color(color)
        first = False
    return box


def add_rect(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    *,
    fill: str = COLORS["white"],
    line: str = COLORS["line"],
    radius: bool = False,
) -> Any:
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = hex_color(fill)
    shape.line.color.rgb = hex_color(line)
    shape.line.width = Pt(0.75)
    return shape


def add_label(slide: Any, text: str, x: float, y: float, w: float, fill: str, color: str = COLORS["white"]) -> None:
    add_rect(slide, x, y, w, 0.36, fill=fill, line=fill, radius=True)
    add_text(slide, text, x + 0.03, y + 0.055, w - 0.06, 0.24, font_size=10, color=color, bold=True, align=PP_ALIGN.CENTER)


def add_title(slide: Any, title: str, subtitle: str = "") -> None:
    add_text(slide, title, 0.65, 0.35, 11.6, 0.55, font_size=24, bold=True)
    if subtitle:
        add_text(slide, subtitle, 0.66, 0.92, 11.7, 0.32, font_size=10, color=COLORS["muted"])
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(0.65), Inches(1.25), Inches(12.7), Inches(1.25))
    line.line.color.rgb = hex_color(COLORS["line"])
    line.line.width = Pt(0.8)


def add_footer(slide: Any, index: int) -> None:
    add_text(
        slide,
        f"기후시민회의 결론 도출과정 연결성 검증 · {index:02d}",
        0.65,
        7.12,
        8.5,
        0.22,
        font_size=8,
        color=COLORS["muted"],
    )


def add_flow_step(slide: Any, number: str, title: str, body: str, x: float, y: float, color: str) -> None:
    add_rect(slide, x, y, 2.35, 1.25, fill=COLORS["white"], line=color, radius=True)
    add_label(slide, number, x + 0.15, y + 0.14, 0.48, color)
    add_text(slide, title, x + 0.72, y + 0.15, 1.45, 0.28, font_size=12, bold=True)
    add_text(slide, body, x + 0.18, y + 0.52, 1.98, 0.54, font_size=9, color=COLORS["muted"])


def add_metric(slide: Any, value: str, label: str, x: float, y: float, color: str) -> None:
    add_rect(slide, x, y, 2.15, 0.95, fill=COLORS["white"], line=COLORS["line"], radius=True)
    add_text(slide, value, x + 0.1, y + 0.13, 1.95, 0.34, font_size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, label, x + 0.12, y + 0.54, 1.91, 0.22, font_size=8, color=COLORS["muted"], align=PP_ALIGN.CENTER)


def add_table_row(slide: Any, columns: list[str], x: float, y: float, widths: list[float], fills: list[str], *, font_size: int = 8) -> None:
    current_x = x
    for index, column in enumerate(columns):
        add_rect(slide, current_x, y, widths[index], 0.54, fill=fills[index], line=COLORS["line"])
        add_text(slide, column, current_x + 0.04, y + 0.08, widths[index] - 0.08, 0.32, font_size=font_size, color=COLORS["ink"], bold=index == 0)
        current_x += widths[index]


def build_storyboard(evidence: dict[str, Any]) -> str:
    coverage = evidence["coverage"]
    regulation = evidence["regulation_decisions"]
    slots = evidence["agenda_final_slots"]
    lines = [
        "# 운영규정·의제 결론 도출과정 연결성 발표자료 스토리보드",
        "",
        "## 핵심 메시지",
        "",
        evidence["summary"]["recommended_framing"],
        "",
        "즉, 이번 자료는 `끝장토론에서 바로 결론`이라고 설명하지 않는다. 숙의가 선택지와 판단 기준을 만들었고, 투표 및 통합 판단이 최종 발표 결론으로 고정한 과정을 보여준다.",
        "",
        "## 입력데이터 반영 검증 요약",
        "",
        f"- 공개 메뉴 기본값: `{evidence['public_menu']['default']}`",
        f"- 공개 source 수: {len(evidence['public_menu']['source_ids'])}",
        f"- A-only live source 노출: {evidence['public_menu']['has_live_a_t1']}",
        f"- 원본 파일: {coverage['original_files']}개",
        f"- source coverage graph: {coverage['source_coverage_graph']['nodes']} nodes / {coverage['source_coverage_graph']['edges']} edges",
        f"- workshop markdown inputs: {coverage['workflow_markdown_files']}개",
        f"- ready sessions: {coverage['ready_sessions']}개",
        f"- partial/review/gap sessions: {coverage['partial_or_review_sessions']}개 ({', '.join(coverage['partial_or_review_labels'])})",
        "",
        "## 운영규정 연결 구조",
        "",
    ]
    for item in regulation:
        lines.append(f"- {item['decision']}: 토론맥락 {item['existing_context_matches']}건, 투표 `{item['vote']}`, 결론 `{item['result']}`")
    lines.extend(["", "## 의제 3개 연결 구조 - 현재 숙의의제 발표덱 기준", ""])
    for slot in slots:
        decision_bits = []
        for decision in slot["source_decisions"]:
            decision_bits.append(f"{decision['agenda']} ({short_vote(decision)}, 토론맥락 {decision['existing_context_matches']}건)")
        source_text = "; ".join(decision_bits)
        lines.append(f"- {slot['slot']} / {slot['conclusion']} / {slot['status']} / deck {slot['deck_slide']}쪽: {source_text}")
        for step in slot["linkage"]:
            lines.append(f"  - {step}")
    lines.extend(["", "## v6 시나리오와 자료 차이", ""])
    for slot in evidence["scenario_variant_slots"]:
        decision_bits = []
        for decision in slot["source_decisions"]:
            decision_bits.append(f"{decision['agenda']} ({short_vote(decision)}, 토론맥락 {decision['existing_context_matches']}건)")
        source_text = "; ".join(decision_bits) if decision_bits else "현재 repo 증거상 확정명 없음"
        lines.append(f"- {slot['slot']} / {slot['conclusion']} / {slot['status']}: {source_text}")
        for step in slot["linkage"]:
            lines.append(f"  - {step}")
    lines.extend(
        [
            "",
            "## 발표 시 주의 문구",
            "",
            "- 모든 데이터가 완전 반영됐다고 말하지 않는다. 현재 검증상 A-only 주장은 해소됐지만, partial/gap 3건은 별도 표시한다.",
            "- 현재 숙의의제 발표덱은 교육, 시민의식·참여, 자원순환을 각각 별도 의제로 보여준다.",
            "- 의제결과 v6 시나리오는 시민의식·참여를 교육 의제에 통합하고 감축2 새 의제를 미확정 슬롯으로 둔다. 두 자료의 차이를 숨기지 않는다.",
            "- 운영규정은 조별토론·전체논의·전자투표 흐름을 통해 결정되었다고 설명한다.",
        ]
    )
    return "\n".join(lines) + "\n"


def build_pptx(evidence: dict[str, Any], output_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    def new_slide(title: str, subtitle: str = "") -> Any:
        slide = prs.slides.add_slide(blank_layout)
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = hex_color(COLORS["paper"])
        add_title(slide, title, subtitle)
        add_footer(slide, len(prs.slides))
        return slide

    slide = prs.slides.add_slide(blank_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = hex_color("EFF6FF")
    add_label(slide, "검증 발표자료", 0.75, 0.62, 1.7, COLORS["blue"])
    add_text(slide, "운영규정과 의제는\n어떻게 결론이 되었나", 0.75, 1.35, 6.2, 1.35, font_size=32, bold=True)
    add_text(
        slide,
        "입력데이터 반영 범위와 토론→투표→통합→결론의 연결성을 분리해 보여주는 보조 발표자료",
        0.8,
        2.9,
        7.2,
        0.45,
        font_size=14,
        color=COLORS["muted"],
    )
    add_rect(slide, 8.05, 1.0, 4.25, 4.75, fill=COLORS["white"], line="BFDBFE", radius=True)
    add_text(slide, "핵심 답변", 8.35, 1.36, 3.3, 0.35, font_size=18, bold=True, color=COLORS["blue"])
    add_text(
        slide,
        "A조만 부각?  아니다\n모든 데이터 완전 반영?  아니다, gap 3건 표시\n결론 연결?  가능하다, 단 투표/통합 단계까지 보여줘야 한다",
        8.35,
        1.95,
        3.45,
        1.8,
        font_size=17,
        color=COLORS["ink"],
    )
    add_text(
        slide,
        "프레이밍: 숙의가 기준과 선택지를 만들고, 투표와 통합 판단이 결론을 고정했다.",
        8.35,
        4.25,
        3.35,
        0.76,
        font_size=13,
        color=COLORS["teal"],
        bold=True,
    )
    add_footer(slide, 1)

    slide = new_slide("1. 입력데이터 반영 상태", "공개 메뉴와 source coverage를 분리해서 확인")
    coverage = evidence["coverage"]
    add_metric(slide, str(coverage["original_files"]), "원본 파일", 0.85, 1.65, COLORS["blue"])
    add_metric(slide, f"{coverage['source_coverage_graph']['nodes']}", "coverage nodes", 3.25, 1.65, COLORS["teal"])
    add_metric(slide, f"{coverage['source_coverage_graph']['edges']}", "coverage edges", 5.65, 1.65, COLORS["teal"])
    add_metric(slide, str(coverage["ready_sessions"]), "ready sessions", 8.05, 1.65, COLORS["green"])
    add_metric(slide, str(coverage["partial_or_review_sessions"]), "partial/gap", 10.45, 1.65, COLORS["rose"])
    add_rect(slide, 0.85, 3.15, 5.7, 1.45, fill=COLORS["white"], line=COLORS["line"], radius=True)
    add_text(slide, "공개 메뉴", 1.1, 3.35, 1.6, 0.3, font_size=15, bold=True)
    add_text(
        slide,
        f"기본값: {evidence['public_menu']['default']}\n노출 source: {len(evidence['public_menu']['source_ids'])}개\nlive-A_t1 노출: {evidence['public_menu']['has_live_a_t1']}",
        1.1,
        3.82,
        4.8,
        0.55,
        font_size=12,
        color=COLORS["muted"],
    )
    add_rect(slide, 6.9, 3.15, 5.35, 1.45, fill=COLORS["white"], line=COLORS["line"], radius=True)
    add_text(slide, "남은 gap", 7.15, 3.35, 1.7, 0.3, font_size=15, bold=True, color=COLORS["rose"])
    add_text(
        slide,
        ", ".join(coverage["partial_or_review_labels"]),
        7.15,
        3.86,
        4.55,
        0.42,
        font_size=12,
        color=COLORS["muted"],
    )
    add_text(slide, "따라서 `A-only` 위험은 해소됐지만, `완전 반영 완료`라고는 말하지 않는다.", 1.0, 5.4, 11.4, 0.45, font_size=18, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    slide = new_slide("2. 결론 도출과정의 정확한 말하기", "끝장토론 결론이 아니라 숙의-투표-통합의 연결")
    steps = [
        ("1", "입력", "사전 제안, 원본 음성·문서, 조별 토론", COLORS["blue"]),
        ("2", "숙의", "쟁점·조건·후보를 만들고 압축", COLORS["teal"]),
        ("3", "투표", "선택지에 대한 선호와 정족수 확인", COLORS["amber"]),
        ("4", "통합", "겹치는 후보와 분과 구조를 정리", COLORS["indigo"]),
        ("5", "결론", "운영규정과 의제 발표자료로 고정", COLORS["rose"]),
    ]
    for index, (number, title, body, color) in enumerate(steps):
        x = 0.75 + index * 2.5
        add_flow_step(slide, number, title, body, x, 2.15, color)
        if index < len(steps) - 1:
            connector = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x + 2.35), Inches(2.78), Inches(x + 2.48), Inches(2.78))
            connector.line.color.rgb = hex_color(COLORS["line"])
            connector.line.width = Pt(1.2)
    add_text(
        slide,
        "발표 문장: “토론이 결론을 바로 낸 것이 아니라, 토론이 투표 가능한 기준과 선택지를 만들었고 투표와 통합 판단으로 최종 결론이 확정되었습니다.”",
        1.15,
        4.7,
        11.0,
        0.85,
        font_size=18,
        bold=True,
        color=COLORS["ink"],
        align=PP_ALIGN.CENTER,
    )

    slide = new_slide("3. 운영규정: 쟁점에서 규정으로", "조별/전체 논의 뒤 투표로 확정된 항목")
    widths = [3.8, 2.25, 1.55, 1.7, 2.35]
    add_table_row(slide, ["결정", "최종값", "투표", "토론맥락", "표현 방식"], 0.65, 1.55, widths, ["DBEAFE", "DBEAFE", "DBEAFE", "DBEAFE", "DBEAFE"], font_size=9)
    y = 2.12
    for item in evidence["regulation_decisions"][:7]:
        add_table_row(
            slide,
            [
                item["decision"],
                item["result"],
                item["vote"],
                f"{item['existing_context_matches']}건",
                "논의→투표→규정",
            ],
            0.65,
            y,
            widths,
            [COLORS["white"], COLORS["white"], COLORS["white"], COLORS["white"], COLORS["white"]],
            font_size=7,
        )
        y += 0.58
    add_text(slide, "운영규정은 정족수·분과·의제 기준을 투표 결과와 함께 설명할 수 있다.", 0.9, 6.33, 11.6, 0.35, font_size=15, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    slide = new_slide("4. 의제 선정: 690건에서 후보와 분과로", "접수·학습·숙의·평가를 거쳐 발표 의제로 압축")
    add_metric(slide, "690", "사전 접수 의제", 0.9, 1.6, COLORS["blue"])
    add_metric(slide, "8", "평가 후보", 3.35, 1.6, COLORS["teal"])
    add_metric(slide, "18", "기획참여단 응답", 5.8, 1.6, COLORS["indigo"])
    add_metric(slide, "3", "분과/의제 슬롯", 8.25, 1.6, COLORS["amber"])
    add_rect(slide, 0.95, 3.15, 11.35, 2.05, fill=COLORS["white"], line=COLORS["line"], radius=True)
    add_text(slide, "의제 선정 기준", 1.25, 3.45, 2.5, 0.3, font_size=16, bold=True)
    add_text(
        slide,
        "기후위기 직·간접 관련성 / 법령 부합 / 시민 일상 실천 가능성 / 사회적 합의 필요성 / 대안 비교 가능성",
        1.25,
        4.02,
        10.3,
        0.45,
        font_size=15,
        color=COLORS["ink"],
    )
    add_text(slide, "운영규정 제14조 기준이 의제 선정의 필터 역할을 한다.", 1.25, 4.55, 9.8, 0.3, font_size=12, color=COLORS["muted"])

    slide = new_slide("5. 의제 3개: 현재 발표덱 기준", "상위 숙의의제 PPTX는 3개 의제를 별도 슬라이드로 제시")
    slot_colors = [COLORS["teal"], COLORS["indigo"], COLORS["cyan"]]
    for index, slot in enumerate(evidence["agenda_final_slots"]):
        x = 0.75 + index * 4.15
        add_rect(slide, x, 1.55, 3.65, 4.6, fill=COLORS["white"], line=slot_colors[index], radius=True)
        add_label(slide, slot["slot"], x + 0.25, 1.82, 0.95, slot_colors[index])
        add_text(slide, slot["conclusion"], x + 0.25, 2.35, 3.05, 0.78, font_size=16, bold=True, color=COLORS["ink"])
        add_text(slide, f"{slot['status']} · deck {slot['deck_slide']}쪽", x + 0.25, 3.22, 3.1, 0.28, font_size=10, color=slot_colors[index], bold=True)
        body = "\n".join(f"- {line}" for line in slot["linkage"])
        add_text(slide, body, x + 0.25, 3.75, 3.1, 1.45, font_size=8, color=COLORS["muted"])
    add_text(slide, "현재 숙의의제 발표덱 기준으로는 교육, 시민의식·참여, 자원순환이 각각 의제 3개로 보인다.", 1.15, 6.42, 10.9, 0.3, font_size=13, bold=True, color=COLORS["teal"], align=PP_ALIGN.CENTER)

    slide = new_slide("6. 의제별 근거 연결", "투표 근거와 토론 맥락을 같이 보여주기")
    widths = [2.25, 3.35, 2.1, 1.65, 3.2]
    add_table_row(slide, ["의제", "결론 설명", "투표 근거", "토론맥락", "해석"], 0.65, 1.55, widths, ["DCFCE7", "DCFCE7", "DCFCE7", "DCFCE7", "DCFCE7"], font_size=9)
    y = 2.15
    table_rows = [
        [
            "교육",
            "전 생애주기 탄소중립 교육 체계 구축",
            "1위",
            "10건",
            "학습·다음세대·지역교육 논의가 투표 결과로 연결",
        ],
        [
            "시민의식·참여",
            "시민의식 개선 및 참여 활성화",
            "공동 2위",
            "12건",
            "일상 행동과 참여 활성화 논의가 별도 의제로 연결",
        ],
        [
            "자원순환·폐기물",
            "배달 문화와 생활폐기물 감축 통합",
            "공동 2위권",
            "12건",
            "소비·생활양식 논의를 감축1 의제로 연결",
        ],
    ]
    for row in table_rows:
        add_table_row(slide, row, 0.65, y, widths, [COLORS["white"]] * 5, font_size=8)
        y += 0.72
    add_rect(slide, 1.0, 5.2, 11.25, 0.8, fill="FEF3C7", line="F59E0B", radius=True)
    add_text(slide, "권고: 3개 의제는 현재 숙의의제 발표덱 기준으로 말하고, v6 통합 시나리오 차이는 다음 장에서 별도 설명한다.", 1.25, 5.45, 10.8, 0.28, font_size=13, color=COLORS["amber"], bold=True, align=PP_ALIGN.CENTER)

    slide = new_slide("7. 자료 간 차이: v6 통합 시나리오", "교육+시민참여 통합과 감축2 새 슬롯은 별도 caveat")
    variant_colors = [COLORS["teal"], COLORS["cyan"], COLORS["rose"]]
    for index, slot in enumerate(evidence["scenario_variant_slots"]):
        x = 0.75 + index * 4.15
        add_rect(slide, x, 1.55, 3.65, 3.55, fill=COLORS["white"], line=variant_colors[index], radius=True)
        add_label(slide, slot["slot"], x + 0.25, 1.82, 0.95, variant_colors[index])
        add_text(slide, slot["conclusion"], x + 0.25, 2.33, 3.05, 0.58, font_size=14, bold=True, color=COLORS["ink"])
        add_text(slide, slot["status"], x + 0.25, 3.08, 3.1, 0.32, font_size=9, color=variant_colors[index], bold=True)
        body = "\n".join(f"- {line}" for line in slot["linkage"][:2])
        add_text(slide, body, x + 0.25, 3.62, 3.1, 0.72, font_size=8, color=COLORS["muted"])
    add_text(slide, "따라서 발표에서는 `현재 숙의의제 덱 기준 3의제`와 `v6 시나리오의 통합/미확정 슬롯`을 섞어 말하지 않는다.", 1.05, 5.8, 11.1, 0.42, font_size=15, bold=True, color=COLORS["rose"], align=PP_ALIGN.CENTER)

    slide = new_slide("8. 남은 검증 게이트", "완전 반영 완료라고 말하기 전에 닫아야 할 항목")
    gaps = coverage["input_gap_items"]
    for index, gap in enumerate(gaps):
        y = 1.65 + index * 1.25
        add_rect(slide, 0.9, y, 11.4, 0.92, fill=COLORS["white"], line=COLORS["rose"], radius=True)
        add_label(slide, gap["tag"], 1.15, y + 0.24, 1.25, COLORS["rose"])
        add_text(slide, gap["label"], 2.62, y + 0.18, 2.8, 0.26, font_size=13, bold=True)
        add_text(
            slide,
            f"상태: {gap['status']} / 전사: {gap['transcript']} / 그래프: {gap['graphStatus']} / 노드: {gap['graphNodes']}",
            5.65,
            y + 0.21,
            5.85,
            0.24,
            font_size=10,
            color=COLORS["muted"],
        )
    add_text(slide, "이 3건이 남아 있기 때문에 `모든 입력이 완전 반영`이라는 표현은 아직 금물이다.", 1.05, 5.8, 11.0, 0.42, font_size=17, bold=True, color=COLORS["rose"], align=PP_ALIGN.CENTER)

    slide = new_slide("9. 발표용 한 문장", "운영규정과 의제 결과를 같이 묶는 닫는 말")
    add_rect(slide, 1.15, 1.65, 11.0, 3.1, fill=COLORS["white"], line=COLORS["teal"], radius=True)
    add_text(
        slide,
        "“이번 결과는 토론 하나로 바로 결론을 낸 것이 아닙니다. 시민 제안과 조별 숙의가 기준과 후보를 만들었고, 기획참여단 투표가 운영규정과 의제의 방향을 확인했으며, 겹치는 후보는 통합해서 오늘의 결론으로 정리했습니다.”",
        1.75,
        2.25,
        9.75,
        1.25,
        font_size=23,
        bold=True,
        color=COLORS["ink"],
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_text(
        slide,
        "단, 데이터 반영 범위는 정직하게 말한다: A-only는 아니지만, partial/gap 3건은 별도 재확인 대상이다.",
        1.65,
        5.3,
        10.2,
        0.42,
        font_size=15,
        color=COLORS["muted"],
        align=PP_ALIGN.CENTER,
    )

    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)


def build_package(out_dir: Path) -> dict[str, Any]:
    evidence = build_evidence_map()
    evidence_path = out_dir / "linkage-evidence-map.json"
    storyboard_path = out_dir / "linkage-storyboard.md"
    pptx_path = out_dir / "20260704_process_to_conclusion_linkage.pptx"
    report_path = out_dir / "build-report.json"

    write_json(evidence_path, evidence)
    write_text(storyboard_path, build_storyboard(evidence))
    build_pptx(evidence, pptx_path)

    report = {
        "generated_at": now_iso(),
        "outputs": {
            "evidence_map": evidence_path.relative_to(ROOT).as_posix(),
            "storyboard": storyboard_path.relative_to(ROOT).as_posix(),
            "pptx": pptx_path.relative_to(ROOT).as_posix(),
        },
        "counts": {
            "slides": 10,
            "public_sources": len(evidence["public_menu"]["source_ids"]),
            "original_files": evidence["coverage"]["original_files"],
            "coverage_nodes": evidence["coverage"]["source_coverage_graph"]["nodes"],
            "coverage_edges": evidence["coverage"]["source_coverage_graph"]["edges"],
            "regulation_decisions": len(evidence["regulation_decisions"]),
            "agenda_candidates": len(evidence["agenda_candidates"]),
            "agenda_final_slots": len(evidence["agenda_final_slots"]),
            "scenario_variant_slots": len(evidence["scenario_variant_slots"]),
            "partial_or_review_sessions": evidence["coverage"]["partial_or_review_sessions"],
        },
        "fallbacks": [
            "PptxGenJS was not used because the local npm exec path returned ECOMPROMISED and bundled pptxgenjs lacked jszip.",
            "python-pptx 1.0.2 was used to generate the presentation without mutating package dependencies.",
        ],
    }
    write_json(report_path, report)
    return report


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Build the presentation linkage evidence package.")
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=DEFAULT_OUT_DIR,
        help=f"Output directory. Default: {DEFAULT_OUT_DIR}",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    report = build_package(args.out_dir)
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
