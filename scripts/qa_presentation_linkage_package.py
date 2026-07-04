#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
import subprocess
import sys
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

from pptx import Presentation


ROOT = Path(__file__).resolve().parents[1]
DEFAULT_OUT_DIR = ROOT / "evaluation" / "presentation-linkage"


def now_iso() -> str:
    return datetime.now(timezone.utc).isoformat(timespec="seconds").replace("+00:00", "Z")


def read_json(path: Path) -> dict[str, Any]:
    with path.open("r", encoding="utf-8") as handle:
        value = json.load(handle)
    if not isinstance(value, dict):
        raise ValueError(f"Expected object JSON: {path}")
    return value


def write_json(path: Path, value: dict[str, Any]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(value, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def write_text(path: Path, value: str) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(value, encoding="utf-8")


def extract_pptx_text(path: Path) -> list[dict[str, Any]]:
    prs = Presentation(path)
    slides: list[dict[str, Any]] = []
    for slide_index, slide in enumerate(prs.slides, start=1):
        texts: list[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = str(shape.text or "").strip()
            if text:
                texts.append(text)
        slides.append({"slide": slide_index, "texts": texts, "joined": "\n".join(texts)})
    return slides


def run_markitdown(pptx_path: Path, log_path: Path) -> dict[str, Any]:
    command = [sys.executable, "-m", "markitdown", str(pptx_path)]
    completed = subprocess.run(command, capture_output=True, text=True, check=False)
    log = [
        f"command: {' '.join(command)}",
        f"returncode: {completed.returncode}",
        "",
        "[stdout]",
        completed.stdout,
        "",
        "[stderr]",
        completed.stderr,
    ]
    write_text(log_path, "\n".join(log))
    return {
        "command": command,
        "returncode": completed.returncode,
        "available": completed.returncode == 0,
        "log": log_path.relative_to(ROOT).as_posix(),
    }


def find_missing_phrases(text: str, phrases: list[str]) -> list[str]:
    return [phrase for phrase in phrases if phrase not in text]


def build_contact_sheet(slide_export_dir: Path, out_path: Path) -> dict[str, Any]:
    pngs = sorted(slide_export_dir.glob("slide-*.png"))
    if not pngs:
        return {"available": False, "reason": "no exported slide PNG files"}
    try:
        from PIL import Image, ImageDraw
    except ModuleNotFoundError:
        return {"available": False, "reason": "Pillow is not installed"}

    thumb_w = 426
    thumb_h = 240
    label_h = 28
    columns = 3
    rows = (len(pngs) + columns - 1) // columns
    sheet = Image.new("RGB", (columns * thumb_w, rows * (thumb_h + label_h)), "white")
    draw = ImageDraw.Draw(sheet)
    for index, path in enumerate(pngs):
        image = Image.open(path).convert("RGB")
        image.thumbnail((thumb_w, thumb_h), Image.Resampling.LANCZOS)
        x = (index % columns) * thumb_w + (thumb_w - image.width) // 2
        y = (index // columns) * (thumb_h + label_h)
        sheet.paste(image, (x, y))
        draw.text(((index % columns) * thumb_w + 8, y + thumb_h + 6), path.stem, fill=(31, 41, 55))
    out_path.parent.mkdir(parents=True, exist_ok=True)
    sheet.save(out_path)
    return {
        "available": True,
        "slides": len(pngs),
        "contact_sheet": out_path.relative_to(ROOT).as_posix(),
    }


def qa_package(out_dir: Path) -> dict[str, Any]:
    build_report_path = out_dir / "build-report.json"
    evidence_path = out_dir / "linkage-evidence-map.json"
    storyboard_path = out_dir / "linkage-storyboard.md"
    pptx_path = out_dir / "20260704_process_to_conclusion_linkage.pptx"
    subagent_path = out_dir / "subagent-verification.md"
    text_extract_path = out_dir / "pptx-text-extract.txt"
    markitdown_log_path = out_dir / "markitdown-check.log"
    slide_export_dir = out_dir / "slide-export"
    contact_sheet_path = out_dir / "slide-contact-sheet.png"

    required_files = [build_report_path, evidence_path, storyboard_path, pptx_path, subagent_path]
    missing_files = [path.relative_to(ROOT).as_posix() for path in required_files if not path.exists()]
    if missing_files:
        raise FileNotFoundError(f"Missing required files: {missing_files}")

    build_report = read_json(build_report_path)
    evidence = read_json(evidence_path)
    slides = extract_pptx_text(pptx_path)
    extracted_text = "\n\n".join(f"## slide {item['slide']}\n{item['joined']}" for item in slides)
    write_text(text_extract_path, extracted_text + "\n")

    must_include = [
        "A조만 부각?",
        "모든 데이터 완전 반영?",
        "숙의가 기준과 선택지를 만들고",
        "운영규정",
        "의제 3개",
        "시민의식·참여",
        "자료 간 차이",
        "새로운 의제 슬롯",
        "B_t2",
        "토론4통합",
        "음성002",
    ]
    forbidden_unqualified = [
        "모든 입력 완전 반영 완료",
        "후보 8건 모두 최종 선정",
        "끝장토론에서 바로 결론",
    ]
    missing_phrases = find_missing_phrases(extracted_text, must_include)
    forbidden_hits = [phrase for phrase in forbidden_unqualified if phrase in extracted_text]

    markitdown = run_markitdown(pptx_path, markitdown_log_path)
    visual_export = build_contact_sheet(slide_export_dir, contact_sheet_path)
    checks = {
        "required_files_exist": not missing_files,
        "slide_count_matches_report": len(slides) == build_report["counts"]["slides"] == 10,
        "public_menu_not_a_only": evidence["public_menu"]["has_live_a_t1"] is False,
        "partial_gap_count_matches": evidence["coverage"]["partial_or_review_sessions"] == 3,
        "agenda_slots_count_matches": len(evidence["agenda_final_slots"]) == 3,
        "scenario_variant_slots_count_matches": len(evidence["scenario_variant_slots"]) == 3,
        "current_agenda_deck_checked": evidence["agenda_deck_cross_check"]["exists"] is True,
        "required_text_present": not missing_phrases,
        "forbidden_unqualified_claims_absent": not forbidden_hits,
    }
    passed = all(checks.values())

    return {
        "generated_at": now_iso(),
        "passed": passed,
        "checks": checks,
        "counts": {
            "slides": len(slides),
            "required_files": len(required_files),
            "missing_required_files": len(missing_files),
            "pptx_characters_extracted": len(extracted_text),
        },
        "missing_files": missing_files,
        "missing_required_phrases": missing_phrases,
        "forbidden_unqualified_hits": forbidden_hits,
        "outputs": {
            "text_extract": text_extract_path.relative_to(ROOT).as_posix(),
            "markitdown_log": markitdown_log_path.relative_to(ROOT).as_posix(),
            "slide_contact_sheet": contact_sheet_path.relative_to(ROOT).as_posix() if contact_sheet_path.exists() else None,
        },
        "markitdown": markitdown,
        "visual_export": visual_export,
        "notes": [
            "markitdown is attempted because the PPTX workflow expects it. If unavailable, python-pptx text extraction is the fallback QA source.",
            "Visual rendering is represented by slide PNG exports and a contact sheet when PowerPoint COM export is available.",
        ],
    }


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="QA the presentation linkage package.")
    parser.add_argument(
        "--out-dir",
        type=Path,
        default=DEFAULT_OUT_DIR,
        help=f"Output directory. Default: {DEFAULT_OUT_DIR}",
    )
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    report = qa_package(args.out_dir)
    report_path = args.out_dir / "qa-report.json"
    write_json(report_path, report)
    print(json.dumps(report, ensure_ascii=False, indent=2))
    if not report["passed"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
